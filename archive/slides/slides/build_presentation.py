#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
FIGURES = ROOT / "report" / "figures"
OUTPUT = ROOT / "slides" / "project1_memory_perf_study.pptx"

TITLE_COLOR = RGBColor(19, 52, 89)
ACCENT_COLOR = RGBColor(0, 102, 153)
TEXT_COLOR = RGBColor(40, 40, 40)
MUTED_COLOR = RGBColor(90, 90, 90)
BACKGROUND_COLOR = RGBColor(248, 249, 251)


def set_slide_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BACKGROUND_COLOR


def add_header_band(slide) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_COLOR
    shape.line.color.rgb = TITLE_COLOR


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    add_header_band(slide)
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.55), Inches(11.8), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(11.8), Inches(0.4))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.size = Pt(11)
        run.font.color.rgb = MUTED_COLOR


def add_bullets(slide, bullets: list[str], left: float, top: float, width: float, height: float) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(8)
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR


def add_image(slide, image_name: str, left: float, top: float, width: float, height: float) -> None:
    image_path = FIGURES / image_name
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def build_deck() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 1. Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_header_band(slide)
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(12.0), Inches(1.0))
    p = title_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Linux perf Based Memory Access Pattern Analysis"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR

    sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(12.0), Inches(0.8))
    p = sub_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Pointer chasing, STREAM-lite, stride locality, row/column traversal, and tiled matmul"
    run.font.size = Pt(17)
    run.font.color.rgb = MUTED_COLOR

    meta_box = slide.shapes.add_textbox(Inches(0.7), Inches(5.7), Inches(8.5), Inches(0.8))
    p = meta_box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Undergraduate systems project | Linux, C, Python, perf automation"
    run.font.size = Pt(16)
    run.font.color.rgb = ACCENT_COLOR

    # 2. Motivation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "Motivation and Questions")
    add_bullets(
        slide,
        [
            "Goal: observe software-visible memory bottlenecks through small reproducible benchmarks.",
            "Question 1: how do runtime and locality change as the working set grows?",
            "Question 2: which patterns behave as latency-bound versus bandwidth-bound?",
            "Question 3: what simple software changes improve reuse most clearly?",
        ],
        left=0.8,
        top=1.7,
        width=6.0,
        height=4.8,
    )
    callout = slide.shapes.add_textbox(Inches(7.2), Inches(2.0), Inches(5.0), Inches(2.5))
    tf = callout.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Key constraint on this host:\nHardware PMU events were unsupported, so final interpretation relies mainly on runtime plus software perf events."
    run.font.size = Pt(18)
    run.font.color.rgb = TEXT_COLOR

    # 3. Method
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "Benchmark Suite and Workflow")
    add_bullets(
        slide,
        [
            "Benchmarks: pointer_chase, stream_lite, stride_access, row_col_traversal, optional matmul_tiled.",
            "Each benchmark emits one CSV line with problem size, timing, and one benchmark-specific metric.",
            "Automation: collect_perf.sh -> parse_perf.py -> plot_results.py.",
            "Controls: core pinning, warmup, repeated runs, fixed compiler flags.",
        ],
        left=0.8,
        top=1.6,
        width=6.0,
        height=4.7,
    )
    add_bullets(
        slide,
        [
            "Preferred events: cycles, instructions, cache-misses, LLC, L1D, dTLB.",
            "Actual usable events on this host: task-clock, cpu-clock, page-faults, context-switches, cpu-migrations.",
            "Result: timing trends are strong, but cache-level attribution must stay conservative.",
        ],
        left=7.0,
        top=1.9,
        width=5.4,
        height=4.1,
    )

    # 4. Pointer chase + stream
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "Strongest Results: Latency Staircase and Streaming Bandwidth")
    add_image(slide, "latency_staircase.png", left=0.7, top=1.55, width=5.9, height=3.8)
    add_image(slide, "bandwidth_vs_working_set.png", left=6.8, top=1.55, width=5.8, height=3.8)
    add_bullets(
        slide,
        [
            "Pointer chase rose from about 20.2 ns/load at 512 KiB to 110.8 ns/load at 8 MiB and 313.4 ns/load at 32 MiB.",
            "STREAM-lite sustained about 32-33 GB/s for copy and 46-49 GB/s for triad on small to medium working sets.",
            "Interpretation: irregular dependent access is latency-bound, while regular streams expose far more overlap.",
        ],
        left=0.8,
        top=5.55,
        width=11.8,
        height=1.4,
    )

    # 5. Stride + traversal
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "Moderate Results: Stride and Traversal Order")
    add_image(slide, "stride_access_cost_vs_stride.png", left=0.7, top=1.55, width=5.9, height=3.8)
    add_image(slide, "row_vs_column_traversal.png", left=6.8, top=1.55, width=5.8, height=3.8)
    add_bullets(
        slide,
        [
            "Stride effect was small for tiny footprints, but stronger for large ones: at 32 MiB, stride 256 reached about 6.06 ns/access.",
            "Row-major stayed faster than column-major, especially at 1024x1024: about 1.76 ns/element versus 2.83 ns/element.",
            "These trends match locality expectations, but the effect sizes are weaker than pointer chasing on this host.",
        ],
        left=0.8,
        top=5.55,
        width=11.8,
        height=1.4,
    )

    # 6. Tiled matmul
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "Best Optimization Example: Tiled Matrix Multiply")
    add_image(slide, "tile_size_effect.png", left=0.7, top=1.55, width=5.9, height=3.8)
    add_image(slide, "roofline_draft.png", left=6.8, top=1.55, width=5.8, height=3.8)
    add_bullets(
        slide,
        [
            "For 128x128, throughput improved from about 2.17 GFLOP/s at tile 8 to 3.80 GFLOP/s at tile 64.",
            "For 256x256, throughput improved from about 2.20 GFLOP/s at tile 8 to 3.85 GFLOP/s at tile 64.",
            "Takeaway: explicit reuse through blocking produced the clearest software-level gain in the project.",
        ],
        left=0.8,
        top=5.55,
        width=11.8,
        height=1.4,
    )

    # 7. Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "Takeaways, Limitations, and Next Step")
    add_bullets(
        slide,
        [
            "Clear results: pointer-chasing latency staircase, streaming bandwidth plateau, and tiled-matmul reuse benefit.",
            "Weaker but consistent trends: wider strides and column-major traversal generally cost more.",
            "Main limitation: hardware PMU events were unavailable, so cache-level attribution remains indirect.",
            "Next step: rerun the exact same suite on a Linux host with working PMU support for cycles, instructions, and cache-miss counters.",
        ],
        left=0.8,
        top=1.7,
        width=11.6,
        height=4.7,
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    build_deck()
